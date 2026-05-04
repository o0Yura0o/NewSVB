"""Replace slide 33 with a two-column Paper-vs-NSVB-ZH loss comparison."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
from lxml import etree

PPTX = 'c:/Users/neo29/workspace/SVC/NSVB-ZH/VoiceSynthesis.pptx'

prs = Presentation(PPTX)

# ---- style constants ----
TITLE_POS = (Inches(1.85), Inches(0.66))
TITLE_SIZE = Pt(36)
BODY_SIZE = Pt(16)
MONO_SIZE = Pt(13)

HIGHLIGHT_FILL = RGBColor(0xDE, 0xEA, 0xF6)
ACCENT = RGBColor(0x1F, 0x3A, 0x68)
SUBTLE = RGBColor(0x55, 0x55, 0x55)
PAPER_COL = RGBColor(0xF5, 0xEC, 0xE0)    # warm beige header
NSVB_COL = RGBColor(0xDE, 0xEA, 0xF6)     # cool blue header

# ---- Remove all shapes from the existing slide 33 ----
target = prs.slides[32]
spTree = target.shapes._spTree
for sp in list(spTree):
    if sp.tag.endswith('}sp') or sp.tag.endswith('}pic') or \
       sp.tag.endswith('}grpSp') or sp.tag.endswith('}graphicFrame') or \
       sp.tag.endswith('}cxnSp'):
        spTree.remove(sp)


# ---- Helpers ----
def add_textbox(slide, x, y, w, h, paragraphs, size=BODY_SIZE,
                align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(para, str):
            runs = [(para, False, None, None)]
        elif isinstance(para, tuple):
            runs = [(para[0],
                     para[1] if len(para) > 1 else False,
                     para[2] if len(para) > 2 else None,
                     para[3] if len(para) > 3 else None)]
        else:
            runs = [(r[0],
                     r[1] if len(r) > 1 else False,
                     r[2] if len(r) > 2 else None,
                     r[3] if len(r) > 3 else None) for r in para]
        for text, bold, color, override_size in runs:
            r = p.add_run()
            r.text = text
            r.font.size = override_size if override_size else size
            if bold:
                r.font.bold = True
            if color:
                r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill, line_color=RGBColor(0xBB, 0xCC, 0xDD)):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = line_color
    rect.line.width = Pt(0.75)
    rect.text_frame.text = ''
    return rect


# ---- Title ----
tb = target.shapes.add_textbox(TITLE_POS[0], TITLE_POS[1],
                               Inches(11.5), Inches(0.77))
p = tb.text_frame.paragraphs[0]
tb.text_frame.margin_left = Inches(0)
r = p.add_run()
r.text = 'Phase 3 Loss 對照：論文版 vs NSVB-ZH 版'
r.font.size = TITLE_SIZE
r.font.bold = True
r.font.color.rgb = ACCENT

# ---- Subtitle ----
add_textbox(target, 1.85, 1.50, 11.0, 0.45, [
    [('從 paired supervision → unpaired distribution alignment，'
      '拆解補足所需監督', True, ACCENT, Pt(14))]
])

# ---- Column headers ----
add_rect(target, 0.55, 2.00, 6.10, 0.50, PAPER_COL)
add_textbox(target, 0.55, 2.03, 6.10, 0.50, [
    [('論文版 NSVB（Paired, Stage 2）', True, ACCENT)]
], size=Pt(16), align=PP_ALIGN.CENTER)

add_rect(target, 6.75, 2.00, 6.10, 0.50, NSVB_COL)
add_textbox(target, 6.75, 2.03, 6.10, 0.50, [
    [('NSVB-ZH（Unpaired, Phase 3）', True, ACCENT)]
], size=Pt(16), align=PP_ALIGN.CENTER)

# ---- Left column body: paper losses ----
add_textbox(target, 0.55, 2.65, 6.10, 4.20, [
    [('L_map1(M) = −log q_φ(M(z_a) | x_p, c_p)', True, ACCENT, Pt(14))],
    [('    在 x_p 提供的專業 posterior 下評估 M(z_a) 的似然度。', False, None, Pt(13))],
    [('    實作：KL(N(μ_a→p\', σ_a→p\') ‖ N(μ_p, σ_p))', False, SUBTLE, Pt(12))],
    '',
    [('L_map2(M) = ‖x̂ − x_p‖₁ + λ (D(x̂) − 1)²', True, ACCENT, Pt(14))],
    [('    mel L1 重建 + mel adversarial。', False, None, Pt(13))],
    [('    註：實作中 L1 項被關閉（cross_way_no_recon_loss）', False, SUBTLE, Pt(12))],
    '',
    [('前提：', True), ('同人同曲 paired x_p 存在，f 有唯一目標。', False)],
    '',
    [('φ, θ, D 在 Stage 2 凍結，只優化 M。', True, SUBTLE, Pt(13))],
])

# ---- Right column body: NSVB-ZH losses ----
add_textbox(target, 6.75, 2.65, 6.10, 4.20, [
    [('L_adv_z  —  Conditional Adversarial on z', True, ACCENT, Pt(14))],
    [('    D(z, [F0, phoneme]) 讓 f(z_a) 條件分布逼近 z_p 分布', False, None, Pt(13))],
    [('    ← 取代 L_map1 的 paired KL', False, SUBTLE, Pt(12))],
    '',
    [('L_PatchNCE  —  對比學習', True, ACCENT, Pt(14))],
    [('    f(z_a)[t] 要比其他 frame 更像 z_a[t]', False, None, Pt(13))],
    [('    ← 補「個體辨識度」，反 Mode Collapse', False, SUBTLE, Pt(12))],
    '',
    [('L_PPG  —  內容鎖定', True, ACCENT, Pt(14))],
    [('    MSE/KL(PPG(x̂), PPG(x_a))   ← 補 paired 天然的內容一致', False, None, Pt(13))],
    '',
    [('L_spk_GRL  —  音色解耦', True, ACCENT, Pt(14))],
    [('    z 上 speaker classifier + GRL   ← 補 paired 天然的同人保護', False, None, Pt(13))],
    '',
    [('L_mel (anchor 20%)  —  輕量重建', True, ACCENT, Pt(14))],
    [('    ← 取代 L_map2 的 L1 + adv（僅用於 anchor 穩定訓練）', False, SUBTLE, Pt(12))],
])

# ---- Bottom highlight: translation principle ----
add_rect(target, 0.55, 6.10, 12.30, 1.00, HIGHLIGHT_FILL)
add_textbox(target, 0.75, 6.20, 11.9, 0.85, [
    [('核心翻譯關係：', True, ACCENT),
     ('unpaired 失去 x_p 當作 ground truth，把 paired 隱含提供的監督'
      '拆解為四個顯式 loss——分布對齊（adv_z）+ 個體保留（NCE）+ 內容一致（PPG）+ 音色不變（GRL）。', False)]
], size=Pt(13))

# ---- Citation ----
tb = target.shapes.add_textbox(Inches(8.42), Inches(7.12),
                               Inches(4.8), Inches(0.30))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
r = p.add_run()
r.text = 'Paper: Liu et al. ACL 2022, §3.3–3.4  |  NSVB-ZH Design'
r.font.size = Pt(10)
r.font.color.rgb = SUBTLE

prs.save(PPTX)
print(f"Slide 33 rewritten. Total slides: {len(prs.slides)}")
