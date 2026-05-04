import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu

p = Presentation('c:/Users/neo29/workspace/SVC/NSVB-ZH/VoiceSynthesis.pptx')

# Inspect title + body style on slide 13 (text-heavy slide)
for sid in [13, 15, 20, 27]:
    slide = p.slides[sid - 1]
    print(f"\n===== Slide {sid} styles =====")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text[:50]
                f = run.font
                color_val = None
                try:
                    if f.color and f.color.rgb:
                        color_val = str(f.color.rgb)
                except Exception:
                    pass
                print(f"  [{shape.name}] text={text!r} font={f.name!r} size={f.size} bold={f.bold} color={color_val}")
            if not para.runs:
                print(f"  [{shape.name}] (empty paragraph)")

# Also check fill colors of rounded rectangles
print("\n===== Rounded rect fills on slide 13 =====")
slide = p.slides[12]
for shape in slide.shapes:
    if shape.shape_type == 1:  # AUTO_SHAPE
        try:
            fill = shape.fill
            print(f"  {shape.name}: fill_type={fill.type}")
            if fill.type == 1:
                print(f"    fore_color.rgb={fill.fore_color.rgb}")
        except Exception as e:
            print(f"  {shape.name}: err {e}")
