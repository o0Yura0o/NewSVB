"""Append NSVB-ZH discussion slides to VoiceSynthesis.pptx, matching existing style."""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SRC = 'c:/Users/neo29/workspace/SVC/NSVB-ZH/VoiceSynthesis.pptx'
DST = 'c:/Users/neo29/workspace/SVC/NSVB-ZH/VoiceSynthesis.pptx'

prs = Presentation(SRC)

blank_layout = None
for lay in prs.slide_layouts:
    if lay.name == '空白':
        blank_layout = lay
        break
if blank_layout is None:
    blank_layout = prs.slide_layouts[6]

TITLE_POS = (Inches(1.85), Inches(0.66))
TITLE_HEIGHT = Inches(0.77)
TITLE_SIZE = Pt(40)
BODY_SIZE = Pt(18)

HIGHLIGHT_FILL = RGBColor(0xDE, 0xEA, 0xF6)
HIGHLIGHT_FILL_2 = RGBColor(0xFD, 0xF2, 0xE9)
ACCENT = RGBColor(0x1F, 0x3A, 0x68)
SUBTLE = RGBColor(0x55, 0x55, 0x55)


def add_title(slide, text, width=9.5):
    tb = slide.shapes.add_textbox(TITLE_POS[0], TITLE_POS[1],
                                  Inches(width), TITLE_HEIGHT)
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = TITLE_SIZE
    run.font.bold = True
    run.font.color.rgb = ACCENT
    return tb


def add_textbox(slide, x, y, w, h, paragraphs, size=BODY_SIZE,
                align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_top = Inches(0.05)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(para, str):
            runs = [(para, False, None)]
        elif isinstance(para, tuple):
            runs = [(para[0], para[1] if len(para) > 1 else False,
                     para[2] if len(para) > 2 else None)]
        else:
            runs = [(r[0], r[1] if len(r) > 1 else False,
                     r[2] if len(r) > 2 else None) for r in para]
        for text, bold, color in runs:
            r = p.add_run()
            r.text = text
            r.font.size = size
            if bold:
                r.font.bold = True
            if color:
                r.font.color.rgb = color
    return tb


def add_highlight_rect(slide, x, y, w, h, fill=HIGHLIGHT_FILL):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y),
                                  Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
    rect.line.width = Pt(0.75)
    rect.text_frame.text = ''
    return rect


def add_citation(slide, url):
    tb = slide.shapes.add_textbox(Inches(8.42), Inches(7.07),
                                  Inches(4.8), Inches(0.34))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = url
    r.font.size = Pt(11)
    r.font.color.rgb = SUBTLE


# ============================================================
# Slide 28 — NSVB-ZH 重建動機
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_title(s, 'NSVB-ZH：為何改為非對稱訓練？', width=11.5)

add_textbox(s, 1.40, 1.60, 11.0, 0.5,
            [('問題起點', True, ACCENT)], size=Pt(22))
add_textbox(s, 1.85, 2.15, 10.5, 1.5, [
    '原論文 NSVB 使用 PopBuTFy：同一位歌手以「業餘」與「專業」兩種方式唱同一首歌。',
    [('中文語系目前找不到這種 ', False),
     ('「同人同曲 amateur–professional paired」', True, ACCENT),
     (' 資料集。', False)],
    '能取得的中文歌聲資料集僅為單一屬性：專業 studio 錄音 或 業餘 KTV 錄音，兩邊曲目、歌手皆不重疊。',
])

add_highlight_rect(s, 1.23, 3.90, 10.88, 1.35, HIGHLIGHT_FILL)
add_textbox(s, 1.50, 4.05, 10.3, 1.10, [
    [('決策：', True, ACCENT),
     ('放棄尋找中文 paired 資料集，改為 ', False),
     ('非對稱（Unpaired）訓練架構', True, ACCENT), ('。', False)],
    [('理由：', True, ACCENT),
     ('VocalVerse（業餘）+ M4Singer（專業）兩個大規模中文資料集可分別扮演'
      '「業餘域」與「專業域」，曲目不需重疊即可訓練技術映射。', False)],
])

add_textbox(s, 1.40, 5.40, 11.0, 0.5,
            [('帶來的改動', True, ACCENT)], size=Pt(22))
add_textbox(s, 1.85, 5.95, 10.5, 1.3, [
    '① 訓練目標從「逐 frame 點對點重建」改為「條件分布對齊」',
    '② DataLoader 從「同人同曲成對抽樣」改為「雙流獨立抽樣」',
    '③ 新增多項解耦與反坍縮機制（詳見後續投影片）',
])

add_citation(s, 'NSVB-ZH Design Note 2026-04')


# ============================================================
# Slide 29 — 中文資料集選擇
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_title(s, '中文資料集選擇', width=8)

add_textbox(s, 1.85, 1.65, 10.5, 0.55,
            [[('VocalVerse（業餘域）', True, ACCENT)]], size=Pt(20))
add_textbox(s, 1.85, 2.20, 10.5, 1.45, [
    '33 首中文流行曲 × 多位 KTV 使用者，共 929 筆 a cappella 錄音',
    '總時長 52.31 小時（平均每筆約 3.4 分鐘），MOS 由 5 位評審打分',
    [('篩選門檻：', True), ('MOS ≤ 3.0 → 516 筆 ≈ 29.1 小時（與專業域總時長對稱）', False)],
])

add_textbox(s, 1.85, 3.80, 10.5, 0.55,
            [[('M4Singer（專業域）', True, ACCENT)]], size=Pt(20))
add_textbox(s, 1.85, 4.35, 10.5, 1.45, [
    '20 位專業歌手 × 700 首中文流行曲，studio 錄音',
    '總時長 29.77 小時，FastSpeech / DiffSinger 等主流研究常用之專業中文資料集',
    [('與 VocalVerse 曲目不要求重疊', True, ACCENT),
     ('——這正是非對稱訓練的適用前提。', False)],
])

add_highlight_rect(s, 1.23, 5.95, 10.88, 1.20, HIGHLIGHT_FILL_2)
add_textbox(s, 1.50, 6.10, 10.3, 1.00, [
    [('資料規模對稱性：', True, ACCENT),
     ('篩選後業餘域 ≈ 29.1h，專業域 ≈ 29.8h，雙流 batch 取樣機率對稱，避免任一側統計支配。', False)],
])

add_citation(s, 'karl-wang/VocalVerse-dataset  |  M4Singer')


# ============================================================
# Slide 30 — Unpaired Training 理論基礎
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_title(s, 'Unpaired Training 理論基礎', width=9.5)

add_textbox(s, 1.85, 1.55, 10.5, 0.55,
            [[('為什麼「不同曲目、不同歌手」也能學到歌聲美化？', True, ACCENT)]],
            size=Pt(20))

add_textbox(s, 1.85, 2.20, 10.5, 3.40, [
    [('① NSVB 的 CVAE 結構性約束：', True),
     ('Encoder 以 (content PPG, pitch F0, speaker) 為條件，殘差 z 被迫只承載'
      '「前三者之外」的資訊——即技術表現（氣音、顫音、共鳴、音高穩定度）。', False)],
    '',
    [('② 此為結構性質，與訓練目標無關：', True),
     ('即使換成 unpaired 訓練，z 仍保有「技術殘差」的語意。', False)],
    '',
    [('③ 改用分布對齊取代逐點配對：', True),
     ('映射函式 f 的目標從「把 z_a 對到唯一 z_p」改為「讓 f(z_a) 的條件分布'
      '逼近 z_p 的條件分布」，配對關係即可被鬆開。', False)],
])

add_highlight_rect(s, 1.23, 5.75, 10.88, 1.40, HIGHLIGHT_FILL)
add_textbox(s, 1.50, 5.88, 10.3, 1.25, [
    [('既有文獻支撐：', True, ACCENT)],
    '• CycleGAN-VC / StarGAN-VC：unpaired 語音轉換以域對應取代點對點',
    '• AutoVC：以 content/speaker 解耦實現 zero-shot 轉換',
    '• CUT (Park et al. ECCV 2020)：unpaired I2I 以 contrastive 取代 cycle',
])

add_citation(s, 'Park et al. ECCV 2020 (CUT); Liu et al. ACL 2022 (NSVB)')


# ============================================================
# Slide 31 — 架構調整：Frame-level 映射 + 雙流 DataLoader
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_title(s, '架構調整：Frame-level 映射 + 雙流 DataLoader', width=12)

add_textbox(s, 1.85, 1.60, 10.5, 0.5,
            [[('保留原論文的映射粒度設計', True, ACCENT)]], size=Pt(18))
add_textbox(s, 1.85, 2.10, 10.5, 1.75, [
    [('原 NSVB 的 GlobalLatentMap：', True),
     ('Conv1d kernel=1 → 每個 z-frame 獨立做 channel-wise 映射，無時間上下文。', False)],
    [('含義：', True), ('作者認為技術成分是 frame-local 的（該瞬間的氣音/顫音），'
                        '時間結構由 content+pitch 條件吸收。', False)],
    [('對 Unpaired 設計的好處：', True),
     ('分布對齊的最小單位 = z-frame，業餘與專業片段不需等長，'
      '可大膽切成 2–3 秒片段訓練。', False)],
])

add_textbox(s, 1.85, 4.10, 10.5, 0.5,
            [[('雙流獨立抽樣的 Mini-batch', True, ACCENT)]], size=Pt(18))
add_textbox(s, 1.85, 4.60, 10.5, 2.60, [
    [('Amateur stream：', True),
     ('從 VocalVerse（MOS≤3.0）分層抽樣 N 筆片段，按 (pitch_bin, phoneme_density) 分層。', False)],
    [('Pro stream：', True),
     ('從 M4Singer 獨立隨機抽樣 M 筆片段，不強制同曲／同歌手／同音域。', False)],
    [('Anchor batch (~20%)：', True),
     ('每個 batch 保留少量業餘樣本通過 f → Decoder 做重建，'
      '計算 content/speaker preservation loss。', False)],
    [('Curriculum：', True),
     ('Stage A 先偏好音域重疊片段；Stage B 放寬到全域分布；Stage C 加入對抗性難例。', False)],
])

add_citation(s, 'NSVB-ZH Training Pipeline Design')


# ============================================================
# Slide 32 — 關鍵風險與對策
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_title(s, '關鍵風險與對策（Unpaired 特有）', width=10.5)

add_textbox(s, 1.40, 1.50, 11.0, 0.50,
            [[('依嚴重度排序，三個必須處理；兩個可延後或跳過', True, ACCENT)]],
            size=Pt(16))

add_highlight_rect(s, 1.23, 2.05, 10.88, 1.35, HIGHLIGHT_FILL)
add_textbox(s, 1.40, 2.15, 10.5, 1.20, [
    [('【最高優先】域混淆（Domain Confounding）：', True, ACCENT),
     ('VocalVerse（KTV 底噪／殘響）與 M4Singer（studio）環境差異 > 技術差異，'
      'f 會走捷徑學成「高級降噪器」。', False)],
    [('對策：', True),
     ('前處理統一（DeepFilterNet 去噪 + 高通 + RMS norm + 隨機 reverb domain randomization）；'
      '訓 domain classifier 驗證 z 上的域準確率應 ≈ 50%。', False)],
])

add_highlight_rect(s, 1.23, 3.55, 10.88, 1.35, HIGHLIGHT_FILL)
add_textbox(s, 1.40, 3.65, 10.5, 1.20, [
    [('【高】Mode Collapse（Generic Pro）：', True, ACCENT),
     ('MMD 只對齊低階矩，f(z_a) 易坍縮到「該桶的平均專業狀態」，喪失個人聲線。', False)],
    [('對策：', True),
     ('以 PatchNCE 對比學習取代原擬的 L_identity——要求 f(z_a)[t] 在特徵空間中'
      '比其他 frame 更像對應的 z_a[t]，允許大幅技術變換。', False)],
])

add_highlight_rect(s, 1.23, 5.05, 10.88, 1.35, HIGHLIGHT_FILL)
add_textbox(s, 1.40, 5.15, 10.5, 1.20, [
    [('【中】音色洩漏（Timbre Leakage）：', True, ACCENT),
     ('Encoder 無法 100% 解耦，z_a 殘留業餘歌手音色，unpaired 下 f 可能順手轉走。', False)],
    [('對策：', True),
     ('z 上加 speaker adversarial classifier（GRL）；'
      '驗證指標：cos(spk_embed(mel_in), spk_embed(mel_out)) > 0.85。', False)],
])

add_textbox(s, 1.40, 6.60, 11.0, 0.6, [
    [('略過：', True, SUBTLE),
     ('Bucketing 顆粒度（改用 continuous critic 即可）、PPG-MSE 剛性'
      '（PPG 本身已足夠寬容）。', False, SUBTLE)],
], size=Pt(14))

add_citation(s, 'See risk.md in repo root')


# ============================================================
# Slide 33 — Phase 3 Loss 組合
# ============================================================
s = prs.slides.add_slide(blank_layout)
add_title(s, 'Phase 3 訓練目標（Loss 組合）', width=10)

add_textbox(s, 1.85, 1.55, 10.5, 0.6, [
    [('L_total = λ_adv·L_adv + λ_nce·L_PatchNCE + λ_ppg·L_PPG '
      '+ λ_spk·L_spk_GRL + λ_mel·L_mel', True, ACCENT)],
], size=Pt(16))

add_textbox(s, 1.85, 2.35, 10.5, 4.60, [
    [('L_adv（Conditional Adversarial）：', True),
     ('判別器 D(z, [F0, phoneme_embed])，讓 f(z_a) 的條件分布逼近 z_p 的條件分布。'
      '取代原論文 per-frame KL，是 unpaired 下分布對齊的主引擎。', False)],
    '',
    [('L_PatchNCE（對比學習）：', True),
     ('query = proj(f(z_a)[t])，positive = proj(z_a[t])，negatives = 同 batch 其他 frame。'
      'τ=0.07，負樣本 256。保護「變好版還是你」，反 Mode Collapse。', False)],
    '',
    [('L_PPG（內容鎖定）：', True),
     ('MSE / KL between PPG(mel_out) and PPG(mel_a)，確保歌詞內容可懂度不流失。', False)],
    '',
    [('L_spk_GRL（音色解耦）：', True),
     ('Speaker classifier 接在 z 上，梯度反轉，使 z 對 speaker 不可分。', False)],
    '',
    [('L_mel（Anchor reconstruction）：', True),
     ('20% anchor batch 走 f → Decoder → mel，與原 mel 對照，穩定訓練。', False)],
])

add_citation(s, 'λ 初始：adv=1.0, nce=1.0, ppg=0.5, spk=0.1, mel=1.0')


prs.save(DST)
print(f"\nDone. Total slides now: {len(prs.slides)}")
print(f"Saved to: {DST}")
