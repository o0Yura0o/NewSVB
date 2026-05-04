import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Emu

p = Presentation('c:/Users/neo29/workspace/SVC/NSVB-ZH/VoiceSynthesis.pptx')
print(f"Slides: {len(p.slides)}")
print(f"Size: {p.slide_width} x {p.slide_height}  (EMU)")
print(f"       {p.slide_width / 914400:.2f} x {p.slide_height / 914400:.2f} (inch)")
print()

for i, slide in enumerate(p.slides):
    print(f"=== Slide {i+1} ===")
    print(f"  Layout: {slide.slide_layout.name}")
    for j, shape in enumerate(slide.shapes):
        info = f"  [{j}] {shape.shape_type} name={shape.name!r}"
        if shape.has_text_frame:
            txt = shape.text_frame.text.replace('\n', ' | ')[:100]
            info += f" text={txt!r}"
        if hasattr(shape, 'left') and shape.left is not None:
            info += f" pos=({shape.left/914400:.2f},{shape.top/914400:.2f}) size=({shape.width/914400:.2f}x{shape.height/914400:.2f})"
        print(info)
    print()
